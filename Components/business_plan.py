from langchain_core.prompts import PromptTemplate
from langchain_groq import ChatGroq
import os
from pptx import Presentation
from pptx.util import Inches

# Initialize LLM
llm = ChatGroq(
    groq_api_key=os.environ["GROQ_API_KEY"],
    model_name="meta-llama/llama-4-scout-17b-16e-instruct"
)

problem_prompt = PromptTemplate(
    input_variables=["report"],
    template="""
You are a startup consultant. Based on the following market research report, generate 3 clear and concise problem statements that highlight real market pain points:

Market Research Report:
-------------------------
{report}

Each problem should be:
- One to two sentences max
- Focused on gaps or challenges users face
- Actionable and startup-worthy

Format:
1. ...
2. ...
3. ...
"""
)

problem_chain = problem_prompt | llm

def generate_problem_statements(report_text: str) -> list:
    result = problem_chain.invoke({"report": report_text})
    lines = result.content.strip().splitlines()
    return [line.strip()[3:].strip() for line in lines if line.strip().startswith(("1.", "2.", "3."))]

# Solution generator
solution_prompt = PromptTemplate(
    input_variables=["problem"],
    template="""
You are a startup expert. Based on the following problem statement, suggest a clear and innovative solution in 3-4 lines.

Problem:
{problem}

Respond in simple bullet points like:
- ...
- ...
"""
)

solution_chain = solution_prompt | llm

def generate_solution(problem_statement: str) -> str:
    response = solution_chain.invoke({"problem": problem_statement})
    return response.content.strip()

# Generate pitch deck PPT
def create_pitch_deck(startup_name: str, problem: str, solution: str, report: str, unique_angle: str):
    prs = Presentation()
    
    # Slide 1 - Startup Name
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = startup_name
    
    # Slide 2 - Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    slide.placeholders[1].text = problem

    # Slide 3 - Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Solution"
    slide.placeholders[1].text = solution

    # Slide 4 - Market Gap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Gap"
    slide.placeholders[1].text = "\n".join([
        line.strip()
        for line in report.splitlines()
        if "market" in line.lower() or "gap" in line.lower() or "need" in line.lower()
    ][:10])

    # Slide 5 - Competitors
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Competitors"
    slide.placeholders[1].text = "\n".join([
        line.strip()
        for line in report.splitlines()
        if "|" in line and not line.startswith("|------")
    ][:5])

    # Slide 6 - Unique Angle
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Unique Angle"
    slide.placeholders[1].text = unique_angle

        # ✅ Ensure folder exists
    output_dir = "./presentations"
    os.makedirs(output_dir, exist_ok=True)

    # ✅ Save presentation safely
    output_path = os.path.join(output_dir, "pitch_deck.pptx")
    prs.save(output_path)
    return output_path